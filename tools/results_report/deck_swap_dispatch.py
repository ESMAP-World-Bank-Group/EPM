"""Point the dispatch picture of each country Overview slide at a new PNG.

The deck is edited in place by the user, so the picture is located by its box
(the 5.60 x 2.50 in panel under the NDP chart) rather than by shape index, and
only its image relationship is rewritten: position, crop and z-order stay put.

    python deck_swap_dispatch.py --years 2025_2035
"""

import argparse
import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

DECK = Path(r"C:\Users\wb590892\Documents\EPM_Models\blacksea_2026"
            r"\BlackSea_regional_power_trade_followup_August2026_results.pptx")
PNGS = Path(r"C:\Users\wb590892\Documents\EPM_Models\blacksea_2026\Data\results\slides")
SLIDES = {3: "turkiye", 5: "georgia", 7: "azerbaijan", 9: "armenia"}
BOX = (5.60, 2.50)                        # inches, the dispatch panel
TOL = 0.06


def dispatch_pic(slide):
    for sh in slide.shapes:
        if sh.shape_type != 13:
            continue
        w, h = Emu(sh.width).inches, Emu(sh.height).inches
        if abs(w - BOX[0]) < TOL and abs(h - BOX[1]) < TOL:
            return sh
    return None


def main():
    p = argparse.ArgumentParser()
    p.add_argument("--deck", default=str(DECK))
    p.add_argument("--years", default="2025_2035")
    p.add_argument("--out", default=None)
    a = p.parse_args()

    deck = Path(a.deck)
    out = Path(a.out) if a.out else deck
    if out == deck:
        bak = deck.with_name("%s.bak_%s%s" % (deck.stem,
                                              datetime.now().strftime("%Y%m%d_%H%M"),
                                              deck.suffix))
        shutil.copy2(deck, bak)
        print("backup :", bak.name)

    prs = Presentation(str(deck))
    for idx, country in SLIDES.items():
        slide = prs.slides[idx - 1]
        pic = dispatch_pic(slide)
        png = PNGS / ("dispatch_%s_%s.png" % (country, a.years))
        if pic is None or not png.exists():
            print("slide %-2d %-11s SKIPPED (%s)" % (
                idx, country, "no panel" if pic is None else "no png"))
            continue
        part, rId = slide.part.get_or_add_image_part(str(png))
        pic._element.blipFill.blip.set(qn("r:embed"), rId)
        print("slide %-2d %-11s <- %s" % (idx, country, png.name))

    prs.save(str(out))
    print("saved  :", out)


if __name__ == "__main__":
    main()
