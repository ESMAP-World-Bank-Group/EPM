"""Clone the Georgia slide pair (Overview / Flows) for the other countries.

The deck already carries a hand-built pair for Georgia: slide 3 with the NDP
capacity comparison over the hourly dispatch panel, slide 4 with the flow maps
over the trade and seasonal charts, and a Notes box on the right.  This copies
those two slides once per country, keeps the geometry, the crops and the Notes
formatting, swaps in the charts produced by slide_dispatch / slide_trade /
slide_seasonal, and writes the country bullets.

The two pictures that are screenshots of the HTML report (the NDP comparison and
the flow maps) cannot be redrawn from here, so their boxes are left empty with a
marker text to paste over.

    python slides_countries.py                      # writes <deck>_countries.pptx
"""

import argparse
import copy
import io
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parents[2]                       # blacksea_2026
DECK = ROOT / "BlackSea_regional_power_trade_followup_August2026_results.pptx"
SLIDES = ROOT / "Data" / "results" / "slides"

OVERVIEW, FLOWS = 3, 4                       # the Georgia pair, 1-based

# Which picture of the source slide gets which chart.  The pictures the deck
# borrowed from the HTML report have no country-neutral equivalent, so they are
# dropped and the empty box is flagged.
CHART = {
    OVERVIEW: {"Picture 18": "dispatch_%s_2025_2030.png"},
    FLOWS: {"Picture 22": "trade_%s.png", "Picture 24": "seasonal_%s.png"},
}
DROP = {
    OVERVIEW: {"Picture 2": "capacity comparison vs NDP"},
    FLOWS: {"Picture 3": "flow maps 2026 / 2030 / 2035"},
}

# (bold lead, rest of the sentence, arrow line or None)
NOTES = {
    ("Armenia", OVERVIEW): [
        ("Nuclear carries the base until it retires",
         ". 3 TWh a year, flat, about a third of supply. It leaves at the end "
         "of the horizon and gas doubles to fill the gap (2.5 to 4.9 TWh).",
         "Metsamor lifetime and replacement to confirm?"),
        ("No new firm capacity is on offer",
         ". The thermal and hydro fleet is unchanged over the whole horizon: "
         "no gas candidate, no new hydro.",
         "Add candidates?"),
        ("Solar is the only build",
         ". 0.45 to 1.6 GW, 2.7 TWh by 2040, with a first 200 MW battery "
         "alongside it.", None),
        ("Reservoir capacity falls after 2035",
         " (0.97 to 0.56 GW). Retirement, not an economic choice.",
         "Check hydro lifetimes."),
        ("Structural exporter",
         ", 2.4 TWh net every year, split between Georgia and the Iran swap.",
         None),
    ],
    ("Armenia", FLOWS): [
        ("Two links, two different jobs",
         ". The Georgian line runs flat at its 150 MW limit; the Iranian swap "
         "carries everything seasonal.", None),
        ("The Iran swap sends power out in summer and takes it back in winter",
         ". 1.4 TWh out, 0.3 TWh in, identical every year: an exogenous "
         "exchange rather than a dispatch decision.",
         "Keep it fixed?"),
        ("The Georgian link is saturated in every hour",
         ". Armenian nuclear and gas are flat and cheap, and they land exactly "
         "when Georgian run-of-river is at its lowest.",
         "Line barely used in reality. To adjust?"),
        ("The direction reverses in 2040",
         ". With nuclear gone, Armenia buys 0.56 TWh from Georgia and only "
         "sends 0.53 back.", None),
        ("Nothing flows to Azerbaijan or Turkiye",
         " in the baseline: the whole exchange transits through Georgia.",
         "Corridors to test in the scenarios?"),
    ],
    ("Azerbaijan", OVERVIEW): [
        ("Gas does everything, and shrinks anyway",
         ". 26 TWh in 2025 and still 22 in 2040, on a fleet that drops from "
         "5.8 to 3.3 GW as old units reach end of life.", None),
        ("Solar takes the space the gas leaves",
         ". 0.26 to 3.8 GW by 2035, 6.4 GW and 9.4 TWh by 2040. It is the only "
         "large build the model picks.", None),
        ("Wind stays marginal",
         ", 0.29 GW by 2040, against the announced offshore ambition.",
         "Candidates and costs to revisit?"),
        ("Hydro adds 180 MW by 2030",
         " and then stays flat; a 460 MW battery arrives with the 2040 solar.",
         None),
        ("The export surplus disappears",
         ". Net exports fall from 3.3 TWh in 2025 to 0.5 in 2035, and "
         "Azerbaijan ends the horizon as a net buyer.", None),
    ],
    ("Azerbaijan", FLOWS): [
        ("One partner in the baseline: Georgia",
         ". The 1000 MW corridor is loaded at about 38 %, and no other border "
         "is modelled.",
         "Iran and Russia links to add?"),
        ("Supplier turns customer",
         ". Azeri gas covers the Georgian winter in 2025 (3.3 TWh out), then "
         "domestic demand and plant retirements eat the margin: 1.4 TWh in "
         "2035, 2.1 TWh imported in 2040.", None),
        ("Winter stays the export season",
         ". From 2035 Azerbaijan buys through spring and summer, when Georgian "
         "hydro and the new solar are abundant, and still sells in Q4.", None),
        ("The switch is a merit-order effect, not a shortage",
         ". Azeri gas sits just under the Georgian marginal cost in 2025 and "
         "just above it once Georgian solar is in.",
         "Gas price path to confirm?"),
    ],
    ("Turkiye", OVERVIEW): [
        ("Wind and solar dominate the build",
         ". PV 20 to 95 GW, wind 13 to 66 GW by 2040. Nothing else is added at "
         "that scale.", None),
        ("Gas is pushed out, then comes back",
         ". 100 TWh in 2025, 11 TWh in 2035, 30 TWh again in 2040 as demand "
         "catches up with the fleet.", None),
        ("Coal barely moves",
         ", 113 to 121 TWh across the whole horizon. It sits under everything "
         "in the merit order.",
         "Retirement and carbon assumptions to revisit?"),
        ("Nuclear follows the plan",
         ", 3.6 GW in 2030 to 9.6 GW and 73 TWh in 2040. It is imposed, not "
         "chosen.", None),
        ("Storage arrives late",
         ": 9.8 GW of batteries and 2 GW of pumped storage, in 2040 only. "
         "Biomass drops to zero from 2030.",
         "Biomass fuel cost to check?"),
    ],
    ("Turkiye", FLOWS): [
        ("Türkiye is the sink of the region, less so over time",
         ". Net imports of 5.9 TWh in 2025, 3.2 in 2030, 1.0 in 2035.", None),
        ("Georgia is the main gate",
         ". 6.1 TWh imported in 2025 with the 700 MW corridor at its limit "
         "every hour. By 2035 Turkish solar and wind displace the gas that set "
         "the price: 3.9 TWh in, and 1.0 TWh back out.", None),
        ("The EU links turn outward",
         ". Bulgaria and Greece together: imports 1.4 to 0.3 TWh, exports 1.6 "
         "to 2.4 TWh.",
         "CBAM treatment to confirm?"),
        ("Winter flips first",
         ". By 2035 Q1 is a net export quarter while summer still buys: winter "
         "wind and nuclear against the cooling peak.", None),
        ("Internal exchanges dwarf the external ones",
         ". About 50 TWh move between the nine Turkish zones against 5 TWh "
         "across the borders, so both charts show the net position.", None),
    ],
}


def pics(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def clone(prs, src):
    """Copy a slide, images included, and append it at the end of the deck."""
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    tree = new.shapes._spTree
    for sh in src.shapes:
        tree.append(copy.deepcopy(sh._element))
    # The copied r:embed ids point at the source slide's relationships, so every
    # image has to be registered again against the new part.
    for old, sh in zip(pics(src), pics(new)):
        _, rId = new.part.get_or_add_image_part(io.BytesIO(old.image.blob))
        blip = sh._element.find(qn("p:blipFill")).find(qn("a:blip"))
        blip.set(qn("r:embed"), rId)
    return new


def swap_image(slide, pic, path):
    _, rId = slide.part.get_or_add_image_part(str(path))
    blip = pic._element.find(qn("p:blipFill")).find(qn("a:blip"))
    blip.set(qn("r:embed"), rId)


def placeholder(slide, pic, text):
    """Replace a screenshot we cannot redraw with a marker of the same size."""
    left, top, width, height = pic.left, pic.top, pic.width, pic.height
    pic._element.getparent().remove(pic._element)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "[ paste %s ]" % text
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = RGBColor(0xA8, 0xB4, 0xC4)
    return box


def title_shape(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.startswith("Baseline results"):
            return sh
    return None


def notes_shape(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.startswith("Notes"):
            return sh
    return None


def set_title(slide, text):
    sh = title_shape(slide)
    p = sh.text_frame.paragraphs[0]
    runs = p.runs
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def run_templates(shape):
    """Pull a bold run, a plain run, a line break and an arrow run out of the
    Georgia notes, so the copies inherit size, colour and the Wingdings arrow."""
    bold = plain = br = arrow = para = None
    for p in shape.text_frame.paragraphs[1:]:
        if para is None and p._p.find(qn("a:pPr")) is not None:
            para = p._p
        for r in p.runs:
            if r.font.bold and bold is None:
                bold = r._r
            elif not r.font.bold and r._r.find(qn("a:rPr")).find(qn("a:sym")) is not None:
                arrow = r._r
            elif plain is None:
                plain = r._r
        b = p._p.find(qn("a:br"))
        if b is not None and br is None:
            br = b
    return para, bold, plain, br, arrow


def write_notes(slide, bullets):
    sh = notes_shape(slide)
    para, bold, plain, br, arrow = run_templates(sh)
    body = sh.text_frame._txBody
    keep = sh.text_frame.paragraphs[0]._p          # the "Notes" heading
    for p in list(body.findall(qn("a:p"))):
        if p is not keep:
            body.remove(p)

    for lead, rest, question in bullets:
        p = copy.deepcopy(para)
        for child in list(p):
            if child.tag != qn("a:pPr"):
                p.remove(child)
        for tmpl, text in ((bold, lead), (plain, rest)):
            r = copy.deepcopy(tmpl)
            r.find(qn("a:t")).text = text
            p.append(r)
        if question:
            p.append(copy.deepcopy(br))
            r = copy.deepcopy(arrow)
            r.find(qn("a:t")).text = " " + question
            p.append(r)
        body.append(p)


def build(prs, src, kind, country, label):
    new = clone(prs, src)
    set_title(new, "Baseline results – %s – %s"
              % (label, "Overview" if kind == OVERVIEW else "Flows"))
    for pic in list(pics(new)):
        if pic.name in CHART[kind]:
            path = SLIDES / (CHART[kind][pic.name] % country.lower())
            if not path.exists():
                raise SystemExit("missing chart: %s" % path)
            swap_image(new, pic, path)
        elif pic.name in DROP[kind]:
            placeholder(new, pic, "%s %s" % (label, DROP[kind][pic.name]))
    write_notes(new, NOTES[(country, kind)])
    return new


def reorder(prs, after, count):
    """Move the freshly appended slides just behind the Georgia pair."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    moved = ids[len(ids) - count:]
    for i, el in enumerate(moved):
        lst.remove(el)
        lst.insert(after + i, el)


def main():
    p = argparse.ArgumentParser()
    p.add_argument("--deck", default=str(DECK))
    p.add_argument("--out", default=None)
    p.add_argument("--countries", default="Armenia,Azerbaijan,Turkiye")
    a = p.parse_args()

    deck = Path(a.deck)
    prs = Presentation(str(deck))
    src = {k: prs.slides[k - 1] for k in (OVERVIEW, FLOWS)}

    labels = {"Turkiye": "Türkiye"}
    n = 0
    for c in [x.strip() for x in a.countries.split(",")]:
        for kind in (OVERVIEW, FLOWS):
            build(prs, src[kind], kind, c, labels.get(c, c))
            n += 1
    reorder(prs, FLOWS, n)

    out = Path(a.out) if a.out else deck.with_name(deck.stem + "_countries.pptx")
    prs.save(str(out))
    print("%s  (+%d slides)" % (out, n))


if __name__ == "__main__":
    main()
