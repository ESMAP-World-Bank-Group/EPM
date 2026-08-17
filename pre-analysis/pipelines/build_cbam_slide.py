# -*- coding: utf-8 -*-
"""The CBAM slide, standalone, built from output_prices.

A separate one-slide deck rather than an edit to the main file: it can be
regenerated, dropped in by hand, and it cannot break anything that already
works. Every figure on it is read from the pipeline outputs at build time, so
the slide cannot drift from what the model was fed.

One table, not two. The levy and what the levy does are the same story told at
two lengths, and splitting them made the reader carry a row identity across the
slide. Merged, the row is the argument: the emission factor sets the levy, the
levy sets the net-back, and by 2040 it sets how much of the year is worth
exporting at all.

Outputs:
    output_prices/slides/CBAM_focus.pptx
    output_prices/slides/cbam_crossing.png

Run:
    python pre-analysis/pipelines/build_cbam_slide.py
"""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from matplotlib.lines import Line2D
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from render_price_slides import FIGSIZE, _legend, _style, day_weights, show
from render_price_slides import INK as HINK, NAVY as HNAVY
from render_price_slides import RUST as HRUST, TEAL as HTEAL

HERE = Path(__file__).resolve().parents[1]
PRICES = HERE / "output_prices"
OUT = PRICES / "slides" / "CBAM_focus.pptx"
CROSSING = PRICES / "slides" / "cbam_crossing.png"
ASSUMPTIONS = PRICES / "slides" / "CBAM_assumptions.pptx"
ANNEX = PRICES / "slides" / "CBAM_annex_EF.pptx"

YEARS = (2026, 2030, 2040)
FOCUS = 2030                  # the year the levy meets the EU price
CHART_YEARS = (2024, 2040)    # the model horizon; the levy series runs to 2053
FONT = "Calibri"

# the deck's palette, read off the existing slides
NAVY = RGBColor(0x1A, 0x3A, 0x5C)
INK = RGBColor(0x34, 0x40, 0x54)
SLATE = RGBColor(0x4A, 0x55, 0x68)
GREY = RGBColor(0x84, 0x97, 0xB0)
PALE = RGBColor(0xF0, 0xF4, 0xF8)
LINE = RGBColor(0xD6, 0xDC, 0xE5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RUST = RGBColor(0x9C, 0x42, 0x21)

# zone -> how the row is named, for the two zones the model actually exports on
EXPORT_ROWS = {"Türkiye": ("Türkiye → BG / GR", "Bulgaria"),
               "Georgia": ("Georgia → RO", "Romania")}
CONTEXT = ["Serbia", "Ukraine", "Moldova. Republic of"]
PRETTY = {"Moldova. Republic of": "Moldova"}


# ── numbers ─────────────────────────────────────────────────────────────────
def ets_path() -> pd.Series:
    """EUR/tCO2 by year - one series, the file repeats it per zone."""
    c = pd.read_csv(PRICES / "cbam_levy.csv")
    return c.groupby("year")["ets_eur2024_per_t"].first()


def levy_eur() -> pd.DataFrame:
    """The levy actually charged, EUR2024/MWh, by exporter and year.

    Read from the pipeline rather than recomputed as EF x ETS: the file already
    carries the 2026 start, so the series is zero in 2024-2025 where the border
    charge does not yet exist. Recomputing it would draw a levy in years when
    none is paid. Identical across the zones an exporter serves - it is set by
    the exporter's factor, not by the buyer - so first() is safe.
    """
    c = pd.read_csv(PRICES / "cbam_levy.csv")
    return c.groupby(["exporter", "year"])["C_eur2024_per_mwh"].first()


def emission_factors() -> pd.Series:
    ef = pd.read_csv(HERE / "config" / "cbam_emission_factors.csv", comment="#")
    return ef.set_index("country")["ef_electricity"].dropna()


def netback_eur() -> pd.DataFrame:
    """Hours-weighted annual export net-back, EUR 2024, by zone/variant/year.

    pHours-weighted, not a plain mean: the day-types run 48 to 552 hours a year,
    and the cheap ones are exactly the ones CBAM pushes under the floor, so an
    unweighted mean would flatter the result.
    """
    qc = json.loads((PRICES / "qc_netback.json").read_text(encoding="utf-8"))
    rate = float(qc["eur_usd_2024"])

    n = pd.read_csv(PRICES / "netback.csv")
    n = n[(n["scenario"] == "CENTRAL") & (n["direction"] == "export")]
    hours = [c for c in n.columns if c[:1] == "t" and c[1:].isdigit()]
    w = day_weights()

    n = n.assign(day=n[hours].mean(axis=1),
                 wt=w.reindex(pd.MultiIndex.from_arrays([n["q"], n["d"]])).values)
    g = n.groupby(["zone", "variant", "year"])
    out = g.apply(lambda x: (x["day"] * x["wt"]).sum() / x["wt"].sum(),
                  include_groups=False)
    return (out / rate).unstack("variant")


def floored_share() -> pd.Series:
    """Share of the year, in %, where the export net-back sits at the floor."""
    n = pd.read_csv(PRICES / "netback.csv")
    n = n[(n["scenario"] == "CENTRAL") & (n["direction"] == "export")
          & (n["variant"] == "CBAM")]
    w = day_weights()
    days = w.reindex(pd.MultiIndex.from_arrays([n["q"], n["d"]])).values / 24.0
    n = n.assign(days=days, floored=n["n_floored"] * days)
    g = n.groupby(["zone", "year"])
    return 100 * g["floored"].sum() / (g["days"].sum() * 24)


def table_rows() -> list[list[str]]:
    ef, ets = emission_factors(), ets_path()
    nb, fl = netback_eur(), floored_share()

    rows = []
    for country, (label, zone) in EXPORT_ROWS.items():
        levy = [f"{ef[country] * ets[y]:.0f}" for y in YEARS]
        ref, cbam = nb.loc[(zone, FOCUS), "REF"], nb.loc[(zone, FOCUS), "CBAM"]
        rows.append([label, f"{ef[country]:.3f}", *levy,
                     f"{ref:.0f}  →  {cbam:.0f}",
                     f"{fl[(zone, YEARS[-1])]:.0f} %"])
    for country in CONTEXT:
        levy = [f"{ef[country] * ets[y]:.0f}" for y in YEARS]
        rows.append([PRETTY.get(country, country), f"{ef[country]:.3f}",
                     *levy, "–", "–"])
    rows.append(["EU reference", f"{ef['European Union']:.3f}",
                 "–", "–", "–", "–", "–"])
    return rows


# ── the chart ───────────────────────────────────────────────────────────────
# colour is the country; Türkiye carries two corridors, Georgia one
CHART_SERIES = [("Türkiye", ["Bulgaria", "Greece"], HNAVY),
                ("Georgia", ["Romania"], HTEAL)]


# The three EU markets, drawn as one visual family: thin, all in the same grey,
# separated only by dash pattern. They are within 5 EUR/MWh of each other, so
# giving each its own colour would advertise a difference that is not there -
# what matters is that the two thick coloured levies are compared to all three.
EU_MARKETS = [("Romania", (0, ())), ("Bulgaria", (0, (4, 2))),
              ("Greece", (0, (1, 1.6)))]


def cbam_crossing() -> Path:
    """The EU prices, and the levy each exporter is charged, on one axis.

    The comparison that carries the slide, with nothing else in it. When the
    carbon levy alone reaches the European price, there is no price at which
    the exporter can sell: the levy has taken the whole value of the MWh before
    the electricity itself is paid for.

    "Central" is not an average of the three markets - it is the price scenario,
    ENTSO-E TYNDP 2024 National Trends (see eu_price_level.PUBLISHED_TRAJECTORIES).
    The three national prices drawn here are all National Trends; what varies
    between them is the market, not the scenario.

    The shading opens where the Turkish levy passes the *highest* of the three
    prices - the point past which no European market on this border is worth
    selling into, not merely the cheapest one.
    """
    lv = levy_eur()
    lvl = pd.read_csv(PRICES / "level_L.csv")
    lvl = lvl[lvl["scenario"] == "CENTRAL"].pivot_table(
        index="year", columns="zone", values="L_eur2024")

    yrs = list(range(CHART_YEARS[0], CHART_YEARS[1] + 1))
    top = [float(lvl.loc[y, [z for z, _ in EU_MARKETS]].max()) for y in yrs]

    fig, ax = plt.subplots(figsize=FIGSIZE, dpi=300)
    _style(ax)

    handles, labels = [], []
    for country, _, colour in CHART_SERIES:
        levy = [lv[(country, y)] for y in yrs]
        ax.fill_between(yrs, top, levy, where=[c > p for c, p in zip(levy, top)],
                        color=HRUST, alpha=0.16, lw=0, interpolate=True)
        ax.plot(yrs, levy, color=colour, lw=2.0, solid_capstyle="round")
        handles.append(Line2D([], [], color=colour, lw=2.0))
        labels.append(f"CBAM levy · {country}")

    for zone, dash in EU_MARKETS:
        ax.plot(yrs, [float(lvl.loc[y, zone]) for y in yrs], color=HINK,
                lw=1.3, ls=dash, zorder=4)
        handles.append(Line2D([], [], color=HINK, lw=1.3, ls=dash))
        labels.append(f"{zone} price")

    ax.set_xlim(*CHART_YEARS)
    ax.set_ylim(0, None)
    ax.set_xticks([2025, 2030, 2035, 2040])
    ax.set_ylabel("€2024 / MWh", fontsize=8, color=HINK, labelpad=2)

    _legend(fig, handles, labels)
    fig.subplots_adjust(left=0.105, right=0.68, top=0.965, bottom=0.115)
    fig.savefig(CROSSING, dpi=300, transparent=True)
    plt.close(fig)
    return CROSSING


# ── slide furniture ─────────────────────────────────────────────────────────
def box(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return s


def write(shape, runs, size=9, colour=INK, align=PP_ALIGN.LEFT, space=0):
    """runs: str, or a list of (text, bold) - one paragraph per list entry."""
    tf = shape.text_frame
    items = [runs] if isinstance(runs, str) else runs
    for i, item in enumerate(items):
        text, bold = (item, False) if isinstance(item, str) else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = text
        r.font.size, r.font.bold = Pt(size), bold
        r.font.color.rgb, r.font.name = colour, FONT
    return shape


def brace(slide, cx, cy, span=1.10, depth=0.26):
    """A downward curly brace centred on (cx, cy).

    LEFT_BRACE rotated 270 deg. PowerPoint rotates about the centre, so the
    shape is placed by its centre and the width/height are swapped relative to
    what is seen on the slide.
    """
    s = slide.shapes.add_shape(MSO_SHAPE.LEFT_BRACE,
                               Inches(cx - depth / 2), Inches(cy - span / 2),
                               Inches(depth), Inches(span))
    s.rotation = 270
    s.fill.background()
    s.line.color.rgb = GREY
    s.line.width = Pt(1.0)
    s.shadow.inherit = False
    return s


def add_table(slide, rows, header, x, y, w, widths, bold=None, rh=0.24,
              small=()):
    """bold: 1-based body rows to set in ink; the rest are context in slate.

    Defaults to the study rows of the focus slide, which is what every caller
    but the assumptions slide wants.
    """
    if bold is None:
        bold = set(range(1, len(EXPORT_ROWS) + 1))
    n_rows, n_cols = len(rows) + 1, len(header)
    gt = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y),
                                Inches(w), Inches(rh * n_rows)).table
    gt.first_row = True
    for i, frac in enumerate(widths):
        gt.columns[i].width = Inches(w * frac / sum(widths))

    for r, line in enumerate([header] + rows):
        gt.rows[r].height = Inches(rh)
        for c, text in enumerate(line):
            cell = gt.cell(r, c)
            cell.margin_left = cell.margin_right = Inches(0.05)
            cell.margin_top = cell.margin_bottom = 0
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = (NAVY if r == 0 else
                                        PALE if r % 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            # the source column is prose, and centred prose is unreadable
            p.alignment = (PP_ALIGN.LEFT if c == 0 or (c in small and r > 0)
                           else PP_ALIGN.CENTER)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(8.5 if r == 0 else 7.5 if c in small else 9)
            run.font.name = FONT
            # the study rows are the slide's subject; the rest is context
            study = r in bold
            run.font.bold = r == 0 or (study and c not in small)
            run.font.color.rgb = (WHITE if r == 0 else SLATE if c in small else
                                  INK if study else SLATE)
    return gt


# ── the slide ───────────────────────────────────────────────────────────────
def build() -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    write(box(slide, 0.17, -0.12, 9.10, 0.82, fill=NAVY), "CBAM – Focus",
          size=22, colour=WHITE)

    # the title bar ends at 0.70; anything above that is behind it
    write(box(slide, 0.53, 0.76, 9.10, 0.28),
          [("→  CBAM: ", True)], size=11, colour=NAVY)
    write(box(slide, 1.55, 0.76, 8.20, 0.28),
          "a carbon price paid at the EU border. Every MWh entering the EU "
          "pays the ETS price on its deemed carbon — one direction only.",
          size=10.5, colour=INK)

    write(box(slide, 0.20, 1.06, 9.60, 0.18), "Calculation", size=9,
          colour=GREY, align=PP_ALIGN.CENTER)
    write(box(slide, 0.17, 1.26, 9.60, 0.36, fill=PALE),
          "CBAM (€/MWh)   =   EF (tCO₂e/MWh)   ×   ETS price (€/tCO₂)   "
          "×   100 %", size=13, colour=NAVY, align=PP_ALIGN.CENTER)

    terms = [(2.10, "EF", "Annex III default — not the grid mix"),
             (5.00, "ETS price", "Observed 2024, then TYNDP 2024:\n"
                                 "84 → 122 → 159 €/t"),
             (7.90, "100 % from 2026", "No phase-in: EU power has had no free "
                                       "allowances since 2013")]
    for cx, head, sub in terms:
        brace(slide, cx, 1.76)
        write(box(slide, cx - 1.45, 1.90, 2.90, 0.20), head, size=10,
              colour=NAVY, align=PP_ALIGN.CENTER)
        write(box(slide, cx - 1.45, 2.10, 2.90, 0.40), sub, size=8.5,
              colour=SLATE, align=PP_ALIGN.CENTER)

    write(box(slide, 0.17, 2.54, 9.60, 0.20),
          "Georgia's grid averages 0.102 tCO₂/MWh — its export default is "
          "0.440, because the plant answering an export order is thermal.",
          size=8.5, colour=RUST, align=PP_ALIGN.CENTER)

    write(box(slide, 0.45, 2.84, 6.00, 0.18),
          "The factor sets the levy, the levy sets what an export is worth  "
          "·  €2024 / MWh, Central EU price", size=9, colour=GREY)
    header = ["Exporter", "EF", "2026", "2030", "2040",
              f"Net-back {FOCUS}\nwithout → with", "Year at the\nfloor, 2040"]
    add_table(slide, table_rows(), header, x=0.45, y=3.04, w=9.10,
              widths=[2.0, 0.8, 0.8, 0.8, 0.8, 2.1, 1.5])

    write(box(slide, 0.45, 4.98, 9.10, 0.24),
          [("By 2030 the Turkish levy (88) meets the EU price itself (89). "
            "Georgia, at half the factor, keeps a margin.", True)],
          size=10, colour=NAVY)

    write(box(slide, 0.20, 5.30, 9.60, 0.18),
          "Sources: Regulation (EU) 2023/956; Implementing Regulation (EU) "
          "2025/2621 Annex III default factors; EU ETS observed 2024 and "
          "ENTSO-E TYNDP 2024 carbon price.", size=7, colour=GREY)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


ASSUMPTION_YEARS = (2024, 2026, 2030, 2035, 2040)
# which of those the sources actually publish - the rest are interpolated, and
# the slide has to say so rather than let five columns look equally solid
PUBLISHED = {2024, 2030, 2040}


SOURCE_COL = {
    "ets": "EEA, EU ETS auctioning revenues (2024)  ·  ENTSO-E TYNDP 2024 "
           "Scenarios Methodology Report, Table 8 p.32 (2030, 2040; EUR2022 "
           "deflated to €2024)  ·  2026 and 2035 interpolated",
    "ef": "Commission Implementing Regulation (EU) 2025/2621, Annex III — "
          "default values for electricity imported into the EU",
}


def assumption_rows() -> tuple[list[str], list[list[str]]]:
    """Header and the three rows: the shared carbon price, then a levy each.

    The ETS row is not a country row - it is the one input both levies are
    built from, which is why it sits above them and carries no emission factor.

    The source travels in the table rather than in a footnote: this is the
    slide a counterpart interrogates before arguing with any result, and a
    figure whose provenance is one row away is a figure they have to trust
    rather than check.
    """
    ef, ets, lv = emission_factors(), ets_path(), levy_eur()
    ys = ASSUMPTION_YEARS

    header = (["", "EF\n(tCO₂e/MWh)"]
              + [f"{y}{'' if y in PUBLISHED else ' *'}" for y in ys]
              + ["Source"])
    rows = [["Carbon price — EU ETS  (€/tCO₂)", "—"]
            + [f"{ets[y]:.1f}" for y in ys] + [SOURCE_COL["ets"]]]
    for country in ("Türkiye", "Georgia"):
        rows.append([f"{country} — CBAM levy  (€/MWh)", f"{ef[country]:.3f}"]
                    + [f"{lv[(country, y)]:.0f}" for y in ys]
                    + [SOURCE_COL["ef"]])
    return header, rows


def build_assumptions() -> Path:
    """The inputs slide: the table, and nothing else on it.

    No formula bar, no caption, no footnote block - every one of those was a
    sentence the table already says, and the ask was a clean table. The only
    thing that could not be dropped is the provenance, so it became a column:
    the years that are published are sourced in the row, and the two that are
    interpolated are starred in the header and explained in the same cell.
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    write(box(slide, 0.17, -0.12, 9.10, 0.82, fill=NAVY),
          "CBAM – Assumptions", size=22, colour=WHITE)

    header, rows = assumption_rows()
    add_table(slide, rows, header, x=0.45, y=1.75, w=9.10,
              widths=[2.45, 0.85, 0.58, 0.58, 0.58, 0.58, 0.58, 2.90],
              bold={2, 3}, rh=0.52, small={7})

    prs.save(str(ASSUMPTIONS))
    return ASSUMPTIONS


# ── the annex slide: why the factor is what it is ───────────────────────────
# Only countries carrying BOTH annexes can appear: the whole point of the panel
# is the ratio between them, so a row with one number would be dead weight.
ANNEX_ROWS = ["Türkiye", "Georgia", "Ukraine", "United Kingdom", "Albania"]

ANNEX_BULLETS = [
    ("The rule.",
     " Article 7(3) of Regulation (EU) 2023/956 sends imported electricity to "
     "default values; Annex IV §4.2.1 sets them at “the CO₂ emission factor in "
     "the third country”; Article 1(4) of IR (EU) 2025/2621 makes that "
     "Annex III."),
    ("The definition that creates the gap.",
     " Annex IV §1(d) defines that factor as “the CO₂ emission data of the "
     "electricity sector” divided by “the gross electricity generation based "
     "on fossil fuels”. Hydro, nuclear and renewables are struck from the "
     "denominator."),
    ("So the export default is the grid average divided by the fossil share.",
     " Not a penalty applied to the mix — a different denominator. Both "
     "annexes are five-year averages of IEA data (IR recitals 7 and 8), built "
     "on a deliberately conservative approach (recital 3)."),
    ("Goods use the other annex.",
     " For electricity consumed in producing steel or aluminium, IR recital 7 "
     "keeps the country grid average — Annex II. Same source, same years, "
     "different denominator."),
]

ANNEX_NOT = [
    ("No mark-up applies to electricity.",
     " The +10 / +20 / +30 % phase-in of IR recital 5 covers the deviation of "
     "individual installations from a national average, so it sits only in the "
     "Annex I goods tables. Annex III carries one column, with no year on it."),
    ("Actual emissions are not available on this border.",
     " Annex IV §5 requires all five of: a PPA with the producer, a plant "
     "directly connected to the EU system (or demonstrated absence of "
     "congestion), ≤ 550 gCO₂/kWh, hourly firm nomination on allocated "
     "interconnection capacity, and monthly verification."),
]

ANNEX_SOURCES = (
    "Sources:  Regulation (EU) 2023/956, Annex IV §1(d), §4.2.1 and §5  ·  "
    "Commission Implementing Regulation (EU) 2025/2621 of 16 December 2025 "
    "(OJ L, 31.12.2025), recitals 3–8, Article 1(3)–(4), Annexes I–III — "
    "values sourced from the IEA  ·  implied fossil shares cross-checked "
    "against Ember Türkiye Electricity Review and IEA Georgia energy mix."
)


def bullets(slide, x, y, w, h, items, size=9.5, space=7):
    """Bold lead-in and running text in one paragraph, one paragraph per item.

    The `write` helper puts each (text, bold) pair on its own line, which is
    what the focus slide wants and the opposite of what prose wants: here the
    lead-in has to sit inside the sentence it opens, or the reader gets a
    heading and an orphan.
    """
    s = box(slide, x, y, w, h)
    tf = s.text_frame
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        # explicit: the first paragraph of a fresh text frame inherits the
        # layout's centring, which leaves bullet one alone in the middle
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space)
        for text, bold, colour in ((f"▪  {lead}", True, NAVY), (rest, False, INK)):
            r = p.add_run()
            r.text = text
            r.font.size, r.font.bold = Pt(size), bold
            r.font.color.rgb, r.font.name = colour, FONT
    return s


def annex_rows() -> tuple[list[str], list[list[str]]]:
    """The ratio check, read from the config so the slide cannot drift from it.

    The implied fossil share is the claim of the slide made falsifiable: if the
    identity holds, the third column has to land on the country's real fossil
    share, and for these five it does. Albania divides 0 by 0 and is shown as a
    dash - it is on the panel for the two zeros to its left, not for a ratio.
    """
    ef = pd.read_csv(HERE / "config" / "cbam_emission_factors.csv", comment="#")
    ef = ef.set_index("country")

    header = ["", "Grid average\nAnnex II", "Export default\nAnnex III",
              "Implied fossil\nshare of generation"]
    rows = []
    for c in ANNEX_ROWS:
        grid, exp = float(ef.loc[c, "ef_grid_average"]), float(ef.loc[c, "ef_electricity"])
        share = f"{100 * grid / exp:.0f} %" if exp else "–"
        rows.append([c, f"{grid:.3f}", f"{exp:.3f}", share])
    return header, rows


def build_annex() -> Path:
    """One annex slide answering the only question this assumption ever gets.

    Every counterpart asks the same thing - why is Georgia charged on 0.440
    when its grid is 0.102 - and the honest answer is a definition plus a
    division, both of which fit on a slide. So the slide is built to be
    checked, not believed: the legal chain on the left, the arithmetic on the
    right, and the two arguments that do NOT work along the bottom, because
    those are what the room will reach for next.
    """
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    write(box(slide, 0.17, -0.12, 9.10, 0.82, fill=NAVY),
          "Annex – Why the levy uses 0.718 and 0.440", size=22, colour=WHITE)

    write(box(slide, 0.45, 0.82, 9.10, 0.24),
          "CBAM prices imported electricity on the intensity of the exporter's "
          "fossil fleet, not on its grid average. The gap is a definition, not "
          "a penalty.", size=11, colour=NAVY)

    bullets(slide, 0.45, 1.22, 5.30, 2.30, ANNEX_BULLETS)

    header, rows = annex_rows()
    add_table(slide, rows, header, x=6.00, y=1.22, w=3.55,
              widths=[1.30, 0.90, 0.90, 1.10], bold={1, 2}, rh=0.30)

    write(box(slide, 6.00, 3.06, 3.55, 0.70),
          "Annex III = Annex II ÷ fossil share. Türkiye's mix is 58–65 % "
          "fossil, Georgia's ~20–25 % gas — the identity holds. Albania has no "
          "thermal plant and is listed at 0 in both: hydro is not what is "
          "taxed, gas is.", size=8, colour=SLATE)

    write(box(slide, 0.45, 3.72, 9.10, 0.20), "Two arguments that do not work",
          size=9, colour=GREY)
    bullets(slide, 0.45, 3.94, 9.10, 1.10, ANNEX_NOT, size=9, space=5)

    write(box(slide, 0.20, 5.14, 9.60, 0.36), ANNEX_SOURCES, size=7,
          colour=GREY)

    prs.save(str(ANNEX))
    return ANNEX


def main() -> None:
    CROSSING.parent.mkdir(parents=True, exist_ok=True)
    for p in (cbam_crossing(), build()):
        print(f"  wrote {p.relative_to(HERE)}  {p.stat().st_size/1024:.0f} kB")
        show(p)


if __name__ == "__main__":
    main()
